"""
ContextOS ingestors/pptx.py — PowerPoint extraction via python-pptx.

Each slide → Markdown section:
  ## Slide N: {title}
  Body text, bullet points, speaker notes appended.
Tables converted to Markdown tables.
"""
from __future__ import annotations

import time
from pathlib import Path

from contextos.ingestors import register


def _shape_to_text(shape) -> list[str]:
    """Extract text lines from a shape."""
    lines = []
    if not shape.has_text_frame:
        return lines
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level  = para.level
        prefix = "  " * level + "- " if level > 0 else ""
        lines.append(f"{prefix}{text}")
    return lines


def _table_shape_to_md(shape) -> str:
    table = shape.table
    rows  = table.rows
    if not rows:
        return ""
    lines = []
    header = [c.text.strip() for c in rows[0].cells]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("|" + "|".join(["---"] * len(header)) + "|")
    for row in rows[1:]:
        lines.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
    return "\n".join(lines)


@register(".pptx", "python-pptx")
def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(path))

    # Core properties
    try:
        cp    = prs.core_properties
        title = (cp.title or "").strip() or path.stem.replace("-"," ").replace("_"," ").title()
        author = (cp.author or "").strip() or None
    except Exception:
        title  = path.stem.replace("-"," ").replace("_"," ").title()
        author = None

    slide_count = len(prs.slides)

    fm_lines = [
        "---",
        "project: unknown",
        "type: note",
        "status: draft",
        "source: pptx",
        f"updated_at: {time.strftime('%Y-%m-%d')}",
        "tags:",
        "  - pptx",
        "  - imported",
    ]
    if author:
        fm_lines.append(f"owner: {author}")
    fm_lines += [
        f"slides: {slide_count}",
        "---",
        "",
        f"# {title}",
        "",
    ]

    sections = []
    for i, slide in enumerate(prs.slides, 1):
        slide_title = ""
        body_lines: list[str] = []

        # Extract shapes in layout order
        for shape in slide.shapes:
            if shape.has_table:
                body_lines.append(_table_shape_to_md(shape))
                continue

            if not shape.has_text_frame:
                continue

            # Detect title placeholder
            from pptx.enum.shapes import PP_PLACEHOLDER
            try:
                ph_type = shape.placeholder_format.type if shape.is_placeholder else None
            except Exception:
                ph_type = None

            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                slide_title = shape.text.strip()
            else:
                body_lines.extend(_shape_to_text(shape))

        # Speaker notes
        notes_text = ""
        try:
            notes_slide = slide.notes_slide
            notes_text  = notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        heading = f"## Slide {i}: {slide_title}" if slide_title else f"## Slide {i}"
        section = heading + "\n"
        if body_lines:
            section += "\n" + "\n".join(body_lines)
        if notes_text:
            section += f"\n\n*Notes: {notes_text}*"

        sections.append(section)

    return "\n".join(fm_lines) + "\n\n" + "\n\n".join(sections)
