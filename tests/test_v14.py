"""
ContextOS v1.4 tests — hybrid search, ingestors, evaluator, MCP auth.
All offline. No network calls. Binary format tests use in-memory generation.
"""
import json
import time
from pathlib import Path
from typing import Optional
import pytest


# ---------------------------------------------------------------------------
# Hybrid Search — RRF fusion
# ---------------------------------------------------------------------------

class TestHybridSearch:

    def _make_results(self, ids_scores: list[tuple[str, float]], use_distance: bool = True) -> list[dict]:
        """Build mock LanceDB-style result dicts."""
        results = []
        for doc_id, score in ids_scores:
            r = {"id": doc_id, "doc_id": doc_id, "title": f"Doc {doc_id}",
                 "content": f"content about {doc_id}", "type": "note",
                 "domain": "", "project": "test", "filepath": "/tmp/x.md",
                 "heading": "", "status": "draft", "tags": "[]"}
            if use_distance:
                r["_distance"] = 1.0 - score  # convert similarity to L2 distance
            else:
                r["_distance"] = score
            results.append(r)
        return results

    def test_rrf_merge_combines_lists(self):
        from contextos.store import VectorStore
        vs = [{"id": "a", "_distance": 0.1}, {"id": "b", "_distance": 0.2},
              {"id": "c", "_distance": 0.3}]
        bs = [{"id": "c", "_distance": 0.1}, {"id": "a", "_distance": 0.2},
              {"id": "d", "_distance": 0.3}]
        merged = VectorStore._rrf_merge(vs, bs, limit=4, alpha=0.7)
        ids = [r["id"] for r in merged]
        # All present in at least one list should appear
        assert "a" in ids
        assert "c" in ids

    def test_rrf_scores_sum_correctly(self):
        from contextos.store import VectorStore
        k = 60
        vs = [{"id": "x", "_distance": 0.0}]   # rank 1
        bs = [{"id": "x", "_distance": 0.0}]   # rank 1
        merged = VectorStore._rrf_merge(vs, bs, limit=1, alpha=0.7)
        assert len(merged) == 1
        expected = 0.7 * (1/(k+1)) + 0.3 * (1/(k+1))
        assert abs(merged[0]["_rrf_score"] - expected) < 1e-9

    def test_rrf_alpha_0_gives_bm25_weight(self):
        from contextos.store import VectorStore
        k = 60
        vs = [{"id": "a", "_distance": 0.0}]   # rank 1
        bs = [{"id": "b", "_distance": 0.0}]   # rank 1 (different doc)
        # alpha=0 → pure BM25 — b should score higher than a
        merged = VectorStore._rrf_merge(vs, bs, limit=2, alpha=0.0)
        ids = [r["id"] for r in merged]
        # b should be first since it's top of bm25 and alpha=0
        assert ids[0] == "b"

    def test_rrf_alpha_1_gives_vector_weight(self):
        from contextos.store import VectorStore
        k = 60
        vs = [{"id": "a", "_distance": 0.0}]
        bs = [{"id": "b", "_distance": 0.0}]
        merged = VectorStore._rrf_merge(vs, bs, limit=2, alpha=1.0)
        assert merged[0]["id"] == "a"  # a is top of vector

    def test_rrf_deduplicates(self):
        from contextos.store import VectorStore
        vs = [{"id": "a", "_distance": 0.1}, {"id": "b", "_distance": 0.2}]
        bs = [{"id": "a", "_distance": 0.1}, {"id": "c", "_distance": 0.2}]
        merged = VectorStore._rrf_merge(vs, bs, limit=5, alpha=0.7)
        ids = [r["id"] for r in merged]
        assert ids.count("a") == 1  # no duplicates

    def test_rrf_limit_respected(self):
        from contextos.store import VectorStore
        vs = [{"id": str(i), "_distance": i/10} for i in range(10)]
        bs = [{"id": str(i), "_distance": i/10} for i in range(10)]
        merged = VectorStore._rrf_merge(vs, bs, limit=3, alpha=0.7)
        assert len(merged) <= 3


# ---------------------------------------------------------------------------
# BM25 search (unit test the logic without LanceDB)
# ---------------------------------------------------------------------------

class TestBM25Logic:

    def test_bm25_scores_relevant_higher(self):
        from rank_bm25 import BM25Okapi
        # Need enough unique documents for IDF to work
        corpus = [
            "payment retry exponential backoff stripe webhook",
            "booking cancellation refund customer policy window",
            "user authentication JWT token oauth login",
            "payment failure retry logic three attempts backoff",
            "slot availability booking calendar schedule",
            "email notification template confirmation cancellation",
            "database postgres migration schema rollback",
        ]
        tokenised = [doc.lower().split() for doc in corpus]
        bm25 = BM25Okapi(tokenised)
        scores = bm25.get_scores("payment retry".lower().split())
        # Docs 0 and 3 mention "payment" AND "retry" — must outscore doc 1 (booking)
        assert scores[0] > scores[1]
        assert scores[3] > scores[1]

    def test_bm25_empty_query_returns_zeros(self):
        from rank_bm25 import BM25Okapi
        corpus = ["hello world", "foo bar baz"]
        bm25 = BM25Okapi([d.split() for d in corpus])
        scores = bm25.get_scores([])
        assert all(s == 0.0 for s in scores)


# ---------------------------------------------------------------------------
# PDF Ingestor
# ---------------------------------------------------------------------------

class TestPDFIngestor:

    def _create_test_pdf(self, tmp_path: Path, text: str = "Hello World") -> Path:
        """Create a minimal in-memory PDF using pymupdf."""
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), text)
        doc.set_metadata({"title": "Test Document", "author": "ContextOS"})
        pdf_path = tmp_path / "test.pdf"
        doc.save(str(pdf_path))
        doc.close()
        return pdf_path

    def test_pdf_extraction_returns_markdown(self, tmp_path):
        from contextos.ingestors.pdf import extract_pdf
        pdf = self._create_test_pdf(tmp_path, "Architecture overview for payments")
        result = extract_pdf(pdf)
        assert result.startswith("---")
        assert "project:" in result
        assert "type: note" in result
        assert "source: pdf" in result

    def test_pdf_title_from_metadata(self, tmp_path):
        from contextos.ingestors.pdf import extract_pdf
        pdf = self._create_test_pdf(tmp_path)
        result = extract_pdf(pdf)
        assert "Test Document" in result

    def test_pdf_content_extracted(self, tmp_path):
        from contextos.ingestors.pdf import extract_pdf
        pdf = self._create_test_pdf(tmp_path, "This is important architecture content")
        result = extract_pdf(pdf)
        assert "important architecture content" in result

    def test_pdf_page_sections(self, tmp_path):
        from contextos.ingestors.pdf import extract_pdf
        import fitz
        doc = fitz.open()
        for i in range(3):
            page = doc.new_page()
            page.insert_text((50, 50), f"Page {i+1} content here")
        pdf_path = tmp_path / "multi.pdf"
        doc.save(str(pdf_path)); doc.close()
        result = extract_pdf(pdf_path)
        assert "## Page 1" in result

    def test_ingestor_registry_has_pdf(self):
        from contextos.ingestors import can_ingest, supported_extensions
        assert can_ingest(Path("test.pdf"))
        assert ".pdf" in supported_extensions()

    def test_ingest_dispatch_pdf(self, tmp_path):
        from contextos.ingestors import ingest
        pdf = self._create_test_pdf(tmp_path)
        result = ingest(pdf)
        assert result is not None
        assert "---" in result


# ---------------------------------------------------------------------------
# DOCX Ingestor
# ---------------------------------------------------------------------------

class TestDOCXIngestor:

    def _create_test_docx(self, tmp_path: Path) -> Path:
        from docx import Document
        from docx.shared import Pt
        doc = Document()
        doc.core_properties.title  = "Architecture Notes"
        doc.core_properties.author = "Engineering"
        doc.add_heading("System Architecture", level=1)
        doc.add_paragraph("This document describes the backend architecture.")
        doc.add_heading("Components", level=2)
        doc.add_paragraph("The system has three layers.")
        # Add a table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Component"
        table.cell(0, 1).text = "Technology"
        table.cell(1, 0).text = "Database"
        table.cell(1, 1).text = "PostgreSQL"
        path = tmp_path / "arch.docx"
        doc.save(str(path))
        return path

    def test_docx_extraction_returns_markdown(self, tmp_path):
        from contextos.ingestors.docx import extract_docx
        docx = self._create_test_docx(tmp_path)
        result = extract_docx(docx)
        assert result.startswith("---")
        assert "source: docx" in result

    def test_docx_title_extracted(self, tmp_path):
        from contextos.ingestors.docx import extract_docx
        docx = self._create_test_docx(tmp_path)
        result = extract_docx(docx)
        assert "Architecture Notes" in result

    def test_docx_headings_to_markdown(self, tmp_path):
        from contextos.ingestors.docx import extract_docx
        docx = self._create_test_docx(tmp_path)
        result = extract_docx(docx)
        assert "# System Architecture" in result
        assert "## Components" in result

    def test_docx_table_to_markdown(self, tmp_path):
        from contextos.ingestors.docx import extract_docx
        docx = self._create_test_docx(tmp_path)
        result = extract_docx(docx)
        assert "PostgreSQL" in result
        assert "|" in result  # Markdown table

    def test_docx_paragraph_text(self, tmp_path):
        from contextos.ingestors.docx import extract_docx
        docx = self._create_test_docx(tmp_path)
        result = extract_docx(docx)
        assert "backend architecture" in result

    def test_ingestor_registry_has_docx(self):
        from contextos.ingestors import can_ingest
        assert can_ingest(Path("notes.docx"))

    def test_ingest_dispatch_docx(self, tmp_path):
        from contextos.ingestors import ingest
        docx = self._create_test_docx(tmp_path)
        result = ingest(docx)
        assert result is not None


# ---------------------------------------------------------------------------
# PPTX Ingestor
# ---------------------------------------------------------------------------

class TestPPTXIngestor:

    def _create_test_pptx(self, tmp_path: Path) -> Path:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        prs.core_properties.title  = "Q1 Architecture Review"
        prs.core_properties.author = "Platform Team"

        slide_layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Payment Service Design"
        slide.placeholders[1].text = "• Stripe integration\n• Retry logic\n• Idempotency"

        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "Database Schema"
        slide2.placeholders[1].text = "• payments table\n• bookings table"

        path = tmp_path / "review.pptx"
        prs.save(str(path))
        return path

    def test_pptx_extraction_returns_markdown(self, tmp_path):
        from contextos.ingestors.pptx import extract_pptx
        pptx = self._create_test_pptx(tmp_path)
        result = extract_pptx(pptx)
        assert result.startswith("---")
        assert "source: pptx" in result

    def test_pptx_slide_sections(self, tmp_path):
        from contextos.ingestors.pptx import extract_pptx
        pptx = self._create_test_pptx(tmp_path)
        result = extract_pptx(pptx)
        assert "## Slide 1" in result
        assert "## Slide 2" in result

    def test_pptx_slide_title_in_heading(self, tmp_path):
        from contextos.ingestors.pptx import extract_pptx
        pptx = self._create_test_pptx(tmp_path)
        result = extract_pptx(pptx)
        assert "Payment Service Design" in result

    def test_pptx_content_extracted(self, tmp_path):
        from contextos.ingestors.pptx import extract_pptx
        pptx = self._create_test_pptx(tmp_path)
        result = extract_pptx(pptx)
        assert "Stripe" in result or "Retry" in result

    def test_pptx_slide_count_in_frontmatter(self, tmp_path):
        from contextos.ingestors.pptx import extract_pptx
        pptx = self._create_test_pptx(tmp_path)
        result = extract_pptx(pptx)
        assert "slides: 2" in result

    def test_ingestor_registry_has_pptx(self):
        from contextos.ingestors import can_ingest
        assert can_ingest(Path("deck.pptx"))


# ---------------------------------------------------------------------------
# Ingestor pipeline integration (vault.py routing)
# ---------------------------------------------------------------------------

class TestIngestorPipeline:

    def test_unsupported_extension_returns_none(self):
        from contextos.ingestors import ingest
        result = ingest(Path("file.xyz"))
        assert result is None

    def test_can_ingest_false_for_unknown(self):
        from contextos.ingestors import can_ingest
        assert not can_ingest(Path("file.mp3"))
        assert not can_ingest(Path("file.xyz"))

    def test_supported_extensions_list(self):
        from contextos.ingestors import supported_extensions
        exts = supported_extensions()
        assert ".pdf"  in exts
        assert ".docx" in exts
        assert ".pptx" in exts
        assert ".csv"  in exts
        assert ".xlsx" in exts


# ---------------------------------------------------------------------------
# Evaluator
# ---------------------------------------------------------------------------

class TestEvaluator:

    def _make_store_mock(self, results: list[dict]):
        """Mock store that returns fixed results."""
        from unittest.mock import MagicMock
        store = MagicMock()
        store.hybrid_search.return_value = results
        store.search.return_value = results
        return store

    def _make_embedder_mock(self):
        from unittest.mock import MagicMock
        emb = MagicMock()
        emb.embed_query.return_value = [0.1] * 384
        return emb

    def _make_result(self, title: str, score: float) -> dict:
        return {"id": title, "title": title, "type": "domain",
                "content": f"content {title}", "_distance": 1.0 - score,
                "domain": "", "project": "test", "filepath": "/x.md",
                "heading": "", "status": "draft", "tags": "[]"}

    def test_perfect_hit_rate(self):
        from contextos.evaluator import EvalQuestion, run_eval
        questions = [
            EvalQuestion(query="payment retry", expected_title="Payment Model", k=5),
        ]
        store = self._make_store_mock([self._make_result("Payment Model", 0.9)])
        emb   = self._make_embedder_mock()
        summary = run_eval(questions, emb, store)
        assert summary.hit_rate == 1.0
        assert summary.mrr == 1.0
        assert summary.results[0].rank == 1

    def test_miss_lowers_hit_rate(self):
        from contextos.evaluator import EvalQuestion, run_eval
        questions = [
            EvalQuestion(query="payment retry", expected_title="Payment Model", k=5),
        ]
        store = self._make_store_mock([self._make_result("Some Other Doc", 0.9)])
        emb   = self._make_embedder_mock()
        summary = run_eval(questions, emb, store)
        assert summary.hit_rate == 0.0
        assert summary.mrr == 0.0
        assert summary.results[0].rank == 0

    def test_rank_2_gives_mrr_05(self):
        from contextos.evaluator import EvalQuestion, run_eval
        questions = [EvalQuestion(query="q", expected_title="Target", k=5)]
        store = self._make_store_mock([
            self._make_result("Other", 0.95),
            self._make_result("Target", 0.85),
        ])
        summary = run_eval(questions, self._make_embedder_mock(), store)
        assert summary.mrr == pytest.approx(0.5)
        assert summary.results[0].rank == 2

    def test_multi_question_avg(self):
        from contextos.evaluator import EvalQuestion, run_eval
        q1 = EvalQuestion(query="q1", expected_title="Doc A", k=5)
        q2 = EvalQuestion(query="q2", expected_title="Doc B", k=5)

        store = self._make_store_mock([self._make_result("Doc A", 0.9)])
        emb   = self._make_embedder_mock()
        # First call hits, second call uses same mock so also hits
        summary = run_eval([q1, q2], emb, store)
        assert summary.total == 2

    def test_no_results_tracked(self):
        from contextos.evaluator import EvalQuestion, run_eval
        questions = [EvalQuestion(query="q", expected_title="X", k=5)]
        store = self._make_store_mock([])  # no results
        summary = run_eval(questions, self._make_embedder_mock(), store)
        assert summary.no_result_pct == 1.0
        assert summary.hit_rate == 0.0

    def test_load_questions_from_file(self, tmp_path):
        from contextos.evaluator import load_questions
        q_file = tmp_path / "questions.json"
        q_file.write_text(json.dumps([
            {"query": "payment flow", "expected_title": "Payment Model", "project": "test"}
        ]))
        questions = load_questions(q_file)
        assert len(questions) == 1
        assert questions[0].query == "payment flow"
        assert questions[0].project == "test"

    def test_load_questions_missing_file(self, tmp_path):
        from contextos.evaluator import load_questions
        with pytest.raises(FileNotFoundError):
            load_questions(tmp_path / "missing.json")

    def test_save_results(self, tmp_path):
        from contextos.evaluator import EvalQuestion, run_eval, save_results
        questions = [EvalQuestion(query="q", expected_title="Doc", k=3)]
        store = self._make_store_mock([self._make_result("Doc", 0.85)])
        summary = run_eval(questions, self._make_embedder_mock(), store)
        out = tmp_path / "results.json"
        save_results(summary, out)
        assert out.exists()
        data = json.loads(out.read_text())
        assert "summary" in data
        assert "results" in data
        assert data["summary"]["total"] == 1

    def test_summary_as_dict(self):
        from contextos.evaluator import EvalSummary
        s = EvalSummary(total=10, hit_rate=0.8, mrr=0.7,
                        avg_top1_score=0.85, no_result_pct=0.0,
                        avg_latency_ms=45.0)
        d = s.as_dict()
        assert d["hit_rate"] == 0.8
        assert d["mrr"] == 0.7
        assert "avg_top1_score" in d


# ---------------------------------------------------------------------------
# MCP tool scope enforcement
# ---------------------------------------------------------------------------

class TestMCPAuth:

    def test_validate_mcp_token_none_returns_none(self):
        from contextos.mcp_server import _validate_mcp_token
        result = _validate_mcp_token(None)
        assert result is None

    def test_validate_mcp_token_empty_returns_none(self):
        from contextos.mcp_server import _validate_mcp_token
        result = _validate_mcp_token("")
        assert result is None

    def test_tool_scopes_all_require_read(self):
        from contextos.mcp_server import TOOL_SCOPES
        for tool, scope in TOOL_SCOPES.items():
            assert scope == "read", f"Tool {tool} should require 'read' scope"

    def test_all_tools_have_scopes(self):
        from contextos.mcp_server import TOOL_HANDLERS, TOOL_SCOPES
        for tool_name in TOOL_HANDLERS:
            assert tool_name in TOOL_SCOPES, f"Tool {tool_name} missing from TOOL_SCOPES"

    def test_token_scope_read_allows_read(self):
        from contextos.schema import TokenScope, Token
        from datetime import datetime, timezone
        token = Token(
            id="ctx_test", name="test", hash="abc",
            created_at=datetime.now(timezone.utc),
            scope=TokenScope.read,
        )
        assert token.has_scope(TokenScope.read)
        assert not token.has_scope(TokenScope.write)


# ---------------------------------------------------------------------------
# Vault ingestor routing (vault.py)
# ---------------------------------------------------------------------------

def test_scan_vault_picks_up_markdown(tmp_path):
    from contextos.vault import scan_vault
    vault = tmp_path / "vault"
    vault.mkdir()
    (vault / "arch.md").write_text(
        "---\nproject: test\ntype: architecture\nstatus: approved\n---\n# Arch\n"
    )
    docs = scan_vault(vault)
    assert len(docs) == 1
    assert docs[0].type.value == "architecture"


def test_scan_vault_skips_unsupported_ext(tmp_path):
    from contextos.vault import scan_vault
    vault = tmp_path / "vault"
    vault.mkdir()
    (vault / "data.csv").write_text("col1,col2\nval1,val2")
    (vault / "notes.md").write_text(
        "---\nproject: test\ntype: note\nstatus: draft\n---\n# Notes\n"
    )
    docs = scan_vault(vault)
    assert len(docs) == 2  # .md file + .csv file (now supported)
